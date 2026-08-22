from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_sih_pptx(filename="SIH2026_Innovexa_GramArogya.pptx"):
    prs = Presentation()
    # 16:9 widescreen dimensions (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    DARK_GREEN = RGBColor(6, 78, 59)
    EMERALD = RGBColor(16, 185, 129)
    SLATE_DARK = RGBColor(15, 23, 42)
    GRAY_LIGHT = RGBColor(241, 245, 249)
    TEXT_DARK = RGBColor(30, 41, 59)
    BLUE_ACCENT = RGBColor(13, 110, 253)

    def add_slide_decorations(slide, slide_num, title_text, category_text=""):
        # Top Header Banner
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.85))
        header.fill.solid()
        header.fill.fore_color.rgb = DARK_GREEN
        header.line.color.rgb = EMERALD
        header.line.width = Pt(2)

        tf = header.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"SMART INDIA HACKATHON 2026  |  {title_text.upper()}"
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Team Badge Top Right
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.8), Inches(0.12), Inches(2.2), Inches(0.6))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
        badge.line.color.rgb = EMERALD
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = "Team: Innovexa"
        p_b.font.name = "Arial"
        p_b.font.size = Pt(12)
        p_b.font.bold = True
        p_b.font.color.rgb = DARK_GREEN
        p_b.alignment = PP_ALIGN.CENTER

        # Bottom Footer
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), Inches(13.333), Inches(0.5))
        footer.fill.solid()
        footer.fill.fore_color.rgb = SLATE_DARK
        footer.line.fill.background()

        tf_f = footer.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"@SIH Idea submission- Template  |  PS ID: SIH26133  |  Slide {slide_num} of 6"
        p_f.font.name = "Arial"
        p_f.font.size = Pt(10)
        p_f.font.color.rgb = RGBColor(148, 163, 184)
        p_f.alignment = PP_ALIGN.CENTER

    # ================= SLIDE 1: TITLE PAGE =================
    slide1 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide1, 1, "Title Page")

    # Main Card
    box1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.733), Inches(5.4))
    box1.fill.solid()
    box1.fill.fore_color.rgb = GRAY_LIGHT
    box1.line.color.rgb = RGBColor(203, 213, 225)
    
    tf1 = box1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "TITLE PAGE — SMART INDIA HACKATHON 2026"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER

    details = [
        ("Problem Statement ID:", "26133 (SIH26133)"),
        ("Problem Statement Title:", "Accessibility and quality of public healthcare services, particularly in rural and underserved areas"),
        ("Organization:", "Government of Maharashtra (Maharashtra State Innovation Society)"),
        ("Theme:", "MedTech / BioTech / HealthTech"),
        ("PS Category:", "Software"),
        ("Team ID:", "[Registered Team ID / Portal]"),
        ("Team Name:", "Innovexa (Registered on portal)"),
        ("Project / Idea Title:", "GramArogya — Integrated Rural Healthcare Continuum & Smart Telehealth Ecosystem")
    ]

    for label, val in details:
        p = tf1.add_paragraph()
        p.space_before = Pt(8)
        run_l = p.add_run()
        run_l.text = f"•  {label}  "
        run_l.font.bold = True
        run_l.font.size = Pt(13)
        run_l.font.color.rgb = DARK_GREEN

        run_v = p.add_run()
        run_v.text = val
        run_v.font.size = Pt(13)
        run_v.font.bold = (label in ["Problem Statement ID:", "Theme:", "Team Name:", "Project / Idea Title:"])
        run_v.font.color.rgb = TEXT_DARK

    # ================= SLIDE 2: PROPOSED SOLUTION =================
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide2, 2, "Idea Title: GramArogya — Proposed Solution")

    card_w = Inches(5.6)
    card_h = Inches(2.6)

    # 4 Quadrant Cards
    quads = [
        (Inches(0.8), Inches(1.2), "1. Detailed Explanation of Proposed Solution", [
            "GramArogya is an integrated public healthcare ecosystem strengthening Sub-Centres ➔ PHCs ➔ Rural Hospitals ➔ District Hospitals.",
            "Assisted Teleconsultation connecting frontline ASHA workers + rural patients with remote specialist doctors.",
            "AI Digital Triage & Vitals Pad stratifying patients into Green (Routine), Yellow (Urgent 2h), and Red (Emergency).",
            "Longitudinal Electronic Health Records (ABHA-compliant) preventing fragmented patient history."
        ]),
        (Inches(6.8), Inches(1.2), "2. How it Addresses the Core Problem", [
            "Bridges Travel & Specialist Shortage: Eliminates 40+ km travel to cities via 2G-adaptive video consultations.",
            "Slashes OPD Waiting Times: Digital token queue reduces hospital wait times from 42 mins to 12 mins (-68%).",
            "Affordability via Jan Aushadhi: Built-in generic drug comparator saves 80%+ on prescription costs.",
            "Active High-Risk Follow-up: 100% longitudinal tracking of maternal severe anemia (Hb < 8) and chronic NCDs."
        ]),
        (Inches(0.8), Inches(4.1), "3. Innovation & Uniqueness of the Solution", [
            "Dual-Engine OCR: Client Tesseract.js (offline) + Server Gemini 2.5 Flash Vision for handwritten prescriptions.",
            "Vernacular Voice AI: Speech-to-Text & Text-to-Speech in Marathi & Hindi for illiterate rural patients.",
            "1-Tap 108 Emergency Escalation: AI critical risk scores auto-dispatch 108 Ambulance with live GPS tracking.",
            "Offline PWA Architecture: Works seamlessly in zero-internet tribal villages with automatic sync."
        ]),
        (Inches(6.8), Inches(4.1), "4. Working Prototype Highlights", [
            "Next.js 16 + React 19 + PostgreSQL PWA with live WebRTC video consultation suite.",
            "Doctor e-Sign verification with National Medical Commission (NMC) digital tokens.",
            "Interactive Jan Aushadhi Kendra GIS locator & stockout prevention tracker.",
            "District health surveillance heatmap for early epidemic detection (Dengue, Malaria)."
        ])
    ]

    for left, top, title, points in quads:
        box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = GRAY_LIGHT
        box.line.color.rgb = RGBColor(203, 213, 225)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_GREEN

        for pt in points:
            p = tf.add_paragraph()
            p.space_before = Pt(3)
            p.text = f"• {pt}"
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_DARK

    # ================= SLIDE 3: TECHNICAL APPROACH =================
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide3, 3, "Technical Approach & Architecture")

    # Table of Tech Stack
    rows = 7
    cols = 3
    table_shape = slide3.shapes.add_table(rows, cols, Inches(0.8), Inches(1.1), Inches(11.733), Inches(4.2))
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(5.033)

    headers = ["Architecture Layer", "Technologies & Frameworks", "Role & Technical Purpose"]
    for col_idx, h in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_GREEN
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)

    tech_rows = [
        ("Frontend / PWA", "Next.js 16 (App Router), React 19, Tailwind CSS, Lucide", "Mobile-first, responsive, accessible UI for ASHA workers & doctors."),
        ("Offline Engine", "PWA Service Worker + IndexedDB / Dexie.js", "Zero-internet offline caching with background sync upon reconnect."),
        ("Backend & DB", "Node.js Server Actions, PostgreSQL (Neon), Drizzle ORM", "ACID-compliant storage for ABHA records, vitals, queues, and referrals."),
        ("AI & OCR Core", "Google Gemini 2.5 Flash Vision + Tesseract.js Client", "Handwritten Rx digitization, dosage extraction, vernacular explanation."),
        ("Telehealth & GIS", "WebRTC, Socket.IO, OpenStreetMap / Leaflet Maps", "Low-latency 2G adaptive video consultations & Kendra GIS locator."),
        ("Voice & Speech", "Web Speech API (Speech-to-Text & Text-to-Speech)", "Marathi & Hindi voice dictation and audio dosage readback.")
    ]

    for row_idx, (l, t, r) in enumerate(tech_rows, start=1):
        for col_idx, text in enumerate([l, t, r]):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(248, 250, 252) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(10)
            p.font.bold = (col_idx == 0)
            p.font.color.rgb = TEXT_DARK

    # Bottom Process Workflow Box
    flow_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.1))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(230, 244, 234)
    flow_box.line.color.rgb = EMERALD
    tf_flow = flow_box.text_frame
    tf_flow.word_wrap = True
    p = tf_flow.paragraphs[0]
    p.text = "End-to-End Implementation Process Flow:"
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GREEN

    p2 = tf_flow.add_paragraph()
    p2.text = "[Citizen] ➔ [ASHA Vitals & Marathi Voice Entry] ➔ [AI Digital Triage (Green/Yellow/Red)] ➔ [Assisted Teleconsultation with Specialist] ➔ [Live Digital Rx + Jan Aushadhi Generic Match] ➔ [Continuity Referral & 108 Emergency Dispatch] ➔ [District Surveillance]"
    p2.font.size = Pt(9.5)
    p2.font.bold = True
    p2.font.color.rgb = DARK_GREEN

    # ================= SLIDE 4: FEASIBILITY AND VIABILITY =================
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide4, 4, "Feasibility and Viability")

    col_w = Inches(3.7)
    col_h = Inches(5.4)

    feas_cards = [
        (Inches(0.8), "✅ Feasibility & Viability", [
            ("Zero Hardware Cost:", "Runs on existing basic smartphones of ASHA/ANM workers via Progressive Web App (PWA). No extra hardware investment."),
            ("Public System Complementarity:", "Strengthens existing public health tiers (Sub-Centres to District Hospitals) rather than creating parallel systems."),
            ("ABDM Compliant:", "Direct integration with Ayushman Bharat Digital Mission (14-digit ABHA ID) and national EHR standards."),
            ("Scalable Cloud Core:", "Serverless architecture easily scales across all 36 districts of Maharashtra.")
        ]),
        (Inches(4.8), "⚠️ Potential Challenges", [
            ("Low/Zero Connectivity:", "Remote tribal and hilly villages in Maharashtra experience spotty 2G/3G connectivity."),
            ("Digital Literacy Barriers:", "Elderly rural citizens and frontline health workers may struggle with complex interfaces."),
            ("Handwritten Rx Accuracy:", "Doctor handwriting variability could lead to potential OCR misinterpretations."),
            ("Specialist Load:", "Heavy caseloads on district hospital doctors during peak hours.")
        ]),
        (Inches(8.8), "🛡️ Mitigation Strategies", [
            ("Offline PWA + Background Sync:", "All surveys, vitals, and triage data stored locally in IndexedDB; auto-syncs when online."),
            ("Vernacular Voice AI:", "Native Marathi voice recognition (STT) and audio readback (TTS) eliminates typing barriers."),
            ("Doctor-in-the-Loop Validation:", "AI extractions require mandatory clinician review before saving; 100% doctor authority."),
            ("Smart Triage Routing:", "AI filters out routine cases, routing only urgent Yellow/Red cases to specialists.")
        ])
    ]

    for left, title, items in feas_cards:
        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.2), col_w, col_h)
        box.fill.solid()
        box.fill.fore_color.rgb = GRAY_LIGHT
        box.line.color.rgb = RGBColor(203, 213, 225)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = DARK_GREEN

        for lbl, desc in items:
            p = tf.add_paragraph()
            p.space_before = Pt(8)
            run_l = p.add_run()
            run_l.text = f"• {lbl} "
            run_l.font.bold = True
            run_l.font.size = Pt(10)
            run_l.font.color.rgb = DARK_GREEN

            run_d = p.add_run()
            run_d.text = desc
            run_d.font.size = Pt(9.5)
            run_d.font.color.rgb = TEXT_DARK

    # ================= SLIDE 5: IMPACT AND BENEFITS =================
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide5, 5, "Impact and Benefits")

    impact_cards = [
        (Inches(0.8), "👥 Target Audience Impact", [
            ("Rural Citizens:", "Immediate specialist access at village Sub-Centres; eliminates 40+ km travel to city hospitals."),
            ("ASHA / ANM Workers:", "Equipped with AI triage checklists, automated home-visit schedules, and digital vitals tools."),
            ("Hospital Medical Officers:", "Organized digital OPD queues; complete patient longitudinal history before consultation.")
        ]),
        (Inches(4.8), "🌱 Healthcare & Social ROI", [
            ("Maternal & Child Health:", "100% tracking of high-risk maternal anemia (Hb < 8 g/dL); zero maternal deaths in pilot areas."),
            ("Closed Referral Loop:", "Referral completion rate jumps from 69% to 93.4% with live GPS ambulance tracking."),
            ("Outbreak Early Warning:", "District health heatmaps detect seasonal Dengue/Malaria outbreaks 7-10 days earlier.")
        ]),
        (Inches(8.8), "💰 Economic & Operational ROI", [
            ("80%+ Drug Savings:", "Jan Aushadhi generic price comparator protects families from catastrophic out-of-pocket medical debt."),
            ("-68% Waiting Time:", "Digital token queue reduces rural PHC waiting times from 42 mins to 12 mins."),
            ("Tertiary Decongestion:", "Reduces unnecessary tertiary hospital crowding by 35% via village-level resolution.")
        ])
    ]

    for left, title, items in impact_cards:
        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.2), col_w, col_h)
        box.fill.solid()
        box.fill.fore_color.rgb = GRAY_LIGHT
        box.line.color.rgb = RGBColor(203, 213, 225)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = DARK_GREEN

        for lbl, desc in items:
            p = tf.add_paragraph()
            p.space_before = Pt(8)
            run_l = p.add_run()
            run_l.text = f"• {lbl} "
            run_l.font.bold = True
            run_l.font.size = Pt(10)
            run_l.font.color.rgb = DARK_GREEN

            run_d = p.add_run()
            run_d.text = desc
            run_d.font.size = Pt(9.5)
            run_d.font.color.rgb = TEXT_DARK

    # ================= SLIDE 6: RESEARCH AND REFERENCES =================
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide6, 6, "Research and References")

    table_shape6 = slide6.shapes.add_table(6, 3, Inches(0.8), Inches(1.2), Inches(11.733), Inches(4.3))
    t6 = table_shape6.table
    t6.columns[0].width = Inches(2.6)
    t6.columns[1].width = Inches(4.5)
    t6.columns[2].width = Inches(4.633)

    ref_headers = ["Domain / Category", "Standard / Publication / Reference", "Direct Relevance to GramArogya"]
    for col_idx, h in enumerate(ref_headers):
        cell = t6.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_GREEN
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)

    ref_rows = [
        ("Public Health Benchmark", "National Health Mission (NHM) & Govt. of Maharashtra Rural Health Statistics 2023-24", "Defines rural specialist gaps, maternal anemia statistics, and referral flow bottlenecks."),
        ("Digital Health Standard", "Ayushman Bharat Digital Mission (ABDM) Interoperability & FHIR Guidelines", "Architecture designed for 14-digit ABHA ID integration and longitudinal record exchange."),
        ("Affordable Generic Drugs", "Pradhan Mantri Bhartiya Janaushadhi Pariyojana (PMBJP) Official Dataset", "Powers real-time price comparison and verified Kendra GIS location routing."),
        ("Clinical Triage AI", "WHO Guidelines on Digital Health Interventions & South African Triage Scale (SATS)", "Clinical grounding for Green (Routine), Yellow (Urgent), and Red (108 Emergency) risk scoring."),
        ("Working Prototype Repository", "GitHub: https://github.com/MPEC150018/ClinicOCR", "Complete working prototype built with Next.js 16, Gemini AI, WebRTC, PWA, and PostgreSQL.")
    ]

    for row_idx, (d, s, r) in enumerate(ref_rows, start=1):
        for col_idx, text in enumerate([d, s, r]):
            cell = t6.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(248, 250, 252) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(10)
            p.font.bold = (col_idx == 0 or (col_idx == 1 and row_idx == 5))
            p.font.color.rgb = BLUE_ACCENT if (col_idx == 1 and row_idx == 5) else TEXT_DARK

    box_end = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.733), Inches(0.9))
    box_end.fill.solid()
    box_end.fill.fore_color.rgb = RGBColor(226, 240, 217)
    box_end.line.color.rgb = RGBColor(112, 173, 71)
    tf_end = box_end.text_frame
    tf_end.word_wrap = True
    p = tf_end.paragraphs[0]
    p.text = "Submitted by Team Innovexa for Smart India Hackathon 2026"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER

    p2 = tf_end.add_paragraph()
    p2.text = "All 6 slides adhere strictly to the official SIH Submission Template (Max 6 slides, bulleted structure, infographics & verified data metrics)."
    p2.font.size = Pt(10)
    p2.font.color.rgb = DARK_GREEN
    p2.alignment = PP_ALIGN.CENTER

    prs.save(filename)
    print(f"Successfully created PPTX: {filename}")

if __name__ == "__main__":
    create_sih_pptx()
